from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
from pathlib import Path

# ── Colors ──
DARK_BLUE  = RGBColor(0x2C, 0x3E, 0x50)
MID_BLUE   = RGBColor(0x34, 0x98, 0xDB)
PURPLE     = RGBColor(0x9B, 0x59, 0xB6)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
RED        = RGBColor(0xC0, 0x39, 0x2B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
DARK_GRAY  = RGBColor(0x55, 0x55, 0x55)

ROOT    = Path(__file__).resolve().parent.parent
REPORTS = ROOT / "reports"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # fully blank


def add_rect(slide, l, t, w, h, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_image_safe(slide, img_path, l, t, w, h):
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    else:
        add_rect(slide, l, t, w, h, LIGHT_GRAY)
        add_text(slide, f"[Chart: {Path(img_path).name}]",
                 l+0.1, t+h/2-0.2, w-0.2, 0.4,
                 size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 5.5, 13.33, 2.0, MID_BLUE)

add_text(slide, "Bangkok Airbnb", 1, 1.2, 11, 1.2,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Market Intelligence Report", 1, 2.5, 11, 0.8,
         size=28, bold=False, color=RGBColor(0xAE, 0xD6, 0xF1),
         align=PP_ALIGN.CENTER)
add_text(slide, "Expernetic Data Engineer Intern Assessment", 1, 3.4, 11, 0.6,
         size=16, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.CENTER)
add_text(slide, "Inside Airbnb Dataset  |  Bangkok, Thailand  |  July 2026",
         1, 5.7, 11, 0.5,
         size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "31,069 Listings  •  11.3M Calendar Records  •  693K Reviews",
         1, 6.3, 11, 0.5,
         size=13, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide, "Agenda", 0.5, 0.2, 12, 0.8,
         size=28, bold=True, color=WHITE)

items = [
    ("01", "Data Engineering Pipeline",    MID_BLUE,  1.4, "Ingestion · Cleaning · Enrichment · Star Schema"),
    ("02", "Exploratory Data Analysis",    GREEN,     2.5, "Price · Geography · Hosts · Reviews"),
    ("03", "Statistical Analysis",         PURPLE,    3.6, "5 Hypothesis Tests · OLS Regression"),
    ("04", "Machine Learning Models",      ORANGE,    4.7, "Price Prediction · K-Means Clustering"),
    ("05", "AI & NLP Insights",            RED,       5.8, "Sentiment Analysis · Topic Modeling · LLM"),
]

for num, title, color, y, sub in items:
    add_rect(slide, 0.5, y, 0.7, 0.7, color)
    add_text(slide, num, 0.5, y, 0.7, 0.7,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.4, y, 6.0, 0.4,
             size=16, bold=True, color=DARK_BLUE)
    add_text(slide, sub, 1.4, y+0.3, 10.0, 0.35,
             size=11, color=DARK_GRAY)


# ════════════════════════════════════════════════════════
# SLIDE 3 — Dataset Overview
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide, "Dataset Overview — Bangkok Inside Airbnb",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

stats = [
    ("31,069", "Active Listings",       MID_BLUE,  0.5,  1.5),
    ("50",     "Neighbourhoods",        GREEN,     3.5,  1.5),
    ("11.3M",  "Calendar Records",      PURPLE,    6.5,  1.5),
    ("693K",   "Guest Reviews",         ORANGE,    9.5,  1.5),
]
for val, label, color, x, y in stats:
    add_rect(slide, x, y, 2.7, 1.5, color)
    add_text(slide, val,   x, y+0.15, 2.7, 0.8,
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y+0.9,  2.7, 0.5,
             size=13, color=WHITE, align=PP_ALIGN.CENTER)

files = [
    "listings.csv.gz — 31K rows × 90 cols — Core listing data",
    "listings.csv — 31K rows × 19 cols — Detailed attributes",
    "calendar.csv.gz — 11.3M rows × 5 cols — Daily availability",
    "reviews.csv.gz — 693K rows × 6 cols — Full review text",
    "neighbourhoods.geojson — 50 polygon boundaries",
]
add_text(slide, "Files Used:", 0.5, 3.3, 12, 0.4,
         size=14, bold=True, color=DARK_BLUE)
for i, f in enumerate(files):
    add_text(slide, f"▸  {f}", 0.7, 3.8 + i*0.55, 12, 0.45,
             size=12, color=DARK_GRAY)


# ════════════════════════════════════════════════════════
# SLIDE 4 — Data Engineering Pipeline
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, MID_BLUE)
add_text(slide, "Data Engineering Pipeline",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

steps = [
    ("INGEST\nsrc/ingest.py",   MID_BLUE,  0.4),
    ("CLEAN\nsrc/clean.py",     GREEN,     3.0),
    ("ENRICH\nsrc/enrich.py",   PURPLE,    5.6),
    ("MODEL\nsrc/model.py",     ORANGE,    8.2),
    ("ANALYSE\nNotebooks",      RED,       10.8),
]
for label, color, x in steps:
    add_rect(slide, x, 1.5, 2.3, 1.3, color)
    add_text(slide, label, x, 1.5, 2.3, 1.3,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if x < 10.8:
        add_text(slide, "→", x+2.3, 1.9, 0.5, 0.5,
                 size=20, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

details = [
    ("✓ Download & extract .gz files\n✓ Skip-if-exists logic\n✓ Full logging", 0.4),
    ("✓ Price parsing (remove $,)\n✓ Date standardization\n✓ Boolean normalization", 3.0),
    ("✓ Join 7 datasets\n✓ Compute occupancy rate\n✓ Host segmentation", 5.6),
    ("✓ Star schema in DuckDB\n✓ 5 tables created\n✓ 6 SQL queries", 8.2),
    ("✓ EDA · Stats · ML\n✓ NLP · Sentiment\n✓ LLM insights", 10.8),
]
for detail, x in details:
    add_text(slide, detail, x, 3.0, 2.3, 1.8, size=10, color=DARK_GRAY)

add_text(slide, "Master Table Output: 31,069 rows × 115 columns → master_listings.csv",
         0.5, 5.0, 12.3, 0.5, size=13, bold=True,
         color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.5, 5.6, 12.3, 0.9, LIGHT_GRAY)
add_text(slide,
         "Star Schema: fact_listings (31K rows)  +  dim_host (7,935)  "
         "+  dim_neighbourhood (50)  +  dim_room_type (4)  +  dim_property_type (79)",
         0.6, 5.7, 12.1, 0.7, size=12, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 5 — Key EDA Findings
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, GREEN)
add_text(slide, "Key EDA Findings — Bangkok Airbnb Market",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

add_image_safe(slide, REPORTS / "fig02_price_by_room_type.png",
               0.4, 1.3, 6.0, 3.5)
add_image_safe(slide, REPORTS / "fig08_superhost_analysis.png",
               6.8, 1.3, 6.0, 3.5)

insights = [
    "▸  Entire homes median ฿1,715/night vs private rooms ฿1,390/night",
    "▸  Superhosts achieve 28.3% occupancy vs 21.9% for regular hosts",
    "▸  Weekday occupancy dominates — Bangkok is a business travel market",
    "▸  Parthum Wan is the most expensive neighbourhood (฿2,652 median)",
]
for i, ins in enumerate(insights):
    add_text(slide, ins, 0.5, 5.0 + i*0.45, 12.3, 0.4,
             size=11, color=DARK_GRAY)


# ════════════════════════════════════════════════════════
# SLIDE 6 — Statistical Analysis
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, PURPLE)
add_text(slide, "Statistical Analysis — Hypothesis Testing",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

hypotheses = [
    ("H1", "Entire home > Private room prices",     "REJECTED ✓", GREEN),
    ("H2", "Superhost > Regular review scores",      "REJECTED ✓", GREEN),
    ("H3", "10+ reviews ≠ fewer reviews prices",     "REJECTED ✓", GREEN),
    ("H4", "Neighbourhood prices differ (ANOVA)",    "REJECTED ✓", GREEN),
    ("H5", "Weekend ≠ Weekday occupancy",            "REJECTED ✓", GREEN),
]
for i, (h, desc, result, color) in enumerate(hypotheses):
    y = 1.4 + i * 0.85
    add_rect(slide, 0.4, y, 0.7, 0.65, PURPLE)
    add_text(slide, h, 0.4, y, 0.7, 0.65,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, 1.3, y+0.1, 8.5, 0.45, size=13, color=DARK_BLUE)
    add_rect(slide, 10.2, y, 2.7, 0.65, color)
    add_text(slide, result, 10.2, y, 2.7, 0.65,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.4, 5.8, 12.3, 1.0, LIGHT_GRAY)
add_text(slide, "All 5 null hypotheses rejected at α=0.05  |  "
         "Effect sizes reported (Cohen's d, eta-squared, rank-biserial r)  |  "
         "OLS R² explains key price variance",
         0.6, 5.9, 12.0, 0.8, size=12, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 7 — ML Models
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, ORANGE)
add_text(slide, "Machine Learning — Price Prediction & Clustering",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

add_image_safe(slide, REPORTS / "fig11_model_comparison.png",
               0.4, 1.3, 7.5, 3.5)
add_image_safe(slide, REPORTS / "fig12_feature_importance.png",
               8.2, 1.3, 4.8, 3.5)

ml_points = [
    "▸  Gradient Boosting outperforms Ridge & Random Forest on all metrics",
    "▸  Top features: accommodates, neighbourhood, occupancy rate",
    "▸  K-Means clustering reveals distinct listing archetypes",
    "▸  Amenity flags (WiFi, AC, Pool) contribute modestly to price",
]
for i, pt in enumerate(ml_points):
    add_text(slide, pt, 0.5, 5.0 + i*0.45, 12.3, 0.4,
             size=11, color=DARK_GRAY)


# ════════════════════════════════════════════════════════
# SLIDE 8 — NLP & AI Insights
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, RED)
add_text(slide, "AI & NLP — Sentiment Analysis & LLM Insights",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

add_image_safe(slide, REPORTS / "fig16_sentiment_distribution.png",
               0.4, 1.3, 6.2, 3.2)
add_image_safe(slide, REPORTS / "fig19_review_volume_trend.png",
               7.0, 1.3, 6.0, 3.2)

nlp_points = [
    "▸  80%+ of reviews are positive — rating inflation confirmed",
    "▸  LDA topics: Location, Cleanliness, Host Communication, Value, Amenities",
    "▸  Negative reviews cluster around noise & cleanliness complaints",
    "▸  COVID-19 caused 90% drop in bookings — strong recovery by 2023",
    "▸  LLM (Claude API) used for automated insight generation",
]
for i, pt in enumerate(nlp_points):
    add_text(slide, pt, 0.5, 4.7 + i*0.45, 12.3, 0.4,
             size=11, color=DARK_GRAY)


# ════════════════════════════════════════════════════════
# SLIDE 9 — Business Recommendations
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide, "Business Recommendations",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

recs = [
    (MID_BLUE,  "For Hosts",
     "Price entire homes ฿1,500–฿2,500 in central neighbourhoods. "
     "Pursue superhost status — it drives 29% higher occupancy. "
     "Offer weekly discounts to capture business travellers."),
    (GREEN,     "For Investors",
     "Target Parthum Wan & Vadhana for premium returns. "
     "Avg annual revenue of ฿529K in top neighbourhoods. "
     "Single-listing hosts achieve highest occupancy rates."),
    (PURPLE,    "For Platform",
     "Use sentiment analysis to flag at-risk listings proactively. "
     "Address rating inflation with verified review scoring. "
     "Surface weekday availability to business travel segments."),
    (ORANGE,    "For Analysts",
     "Gradient Boosting is the recommended pricing model architecture. "
     "Neighbourhood and capacity explain majority of price variance. "
     "Amenity extraction from text adds meaningful predictive signal."),
]
for i, (color, title, text) in enumerate(recs):
    x = 0.4 + (i % 2) * 6.5
    y = 1.5 + (i // 2) * 2.8
    add_rect(slide, x, y, 6.1, 2.5, color)
    add_text(slide, title, x+0.1, y+0.1, 5.9, 0.5,
             size=15, bold=True, color=WHITE)
    add_text(slide, text, x+0.1, y+0.65, 5.8, 1.7,
             size=10.5, color=WHITE)


# ════════════════════════════════════════════════════════
# SLIDE 10 — Tech Stack & Summary
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
add_text(slide, "Tech Stack & Project Summary",
         0.5, 0.2, 12, 0.8, size=24, bold=True, color=WHITE)

stack = [
    ("Data Processing",  "pandas · numpy · DuckDB",          MID_BLUE),
    ("Visualization",    "matplotlib · seaborn · folium",     GREEN),
    ("ML & Statistics",  "scikit-learn · statsmodels · scipy",PURPLE),
    ("NLP & AI",         "NLTK · TextBlob · Claude API",      ORANGE),
    ("Pipeline",         "Python modules · logging · pytest", RED),
    ("Version Control",  "Git · GitHub",                      DARK_BLUE),
]
for i, (cat, tools, color) in enumerate(stack):
    x = 0.4 + (i % 3) * 4.3
    y = 1.5 + (i // 3) * 1.3
    add_rect(slide, x, y, 4.0, 1.1, color)
    add_text(slide, cat,   x+0.1, y+0.05, 3.8, 0.45,
             size=13, bold=True, color=WHITE)
    add_text(slide, tools, x+0.1, y+0.55, 3.8, 0.45,
             size=10, color=WHITE)

summary_items = [
    "✓ 6 sections completed with full depth",
    "✓ 19 visualizations + interactive Folium map",
    "✓ Production-quality pipeline with logging & tests",
    "✓ Star schema with 5 tables in DuckDB",
    "✓ 5 hypothesis tests + OLS regression",
    "✓ 3 ML models + K-Means clustering",
]
add_rect(slide, 0.4, 4.3, 12.3, 2.8, LIGHT_GRAY)
add_text(slide, "Deliverables Summary", 0.6, 4.4, 12, 0.4,
         size=14, bold=True, color=DARK_BLUE)
for i, item in enumerate(summary_items):
    x = 0.7 + (i % 2) * 6.2
    y = 4.9 + (i // 2) * 0.55
    add_text(slide, item, x, y, 6.0, 0.45, size=11, color=DARK_GRAY)


# ════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════
out_path = REPORTS / "Bangkok_Airbnb_Presentation.pptx"
prs.save(str(out_path))
print(f"✓ Presentation saved to {out_path}")
print(f"  Slides: {len(prs.slides)}")